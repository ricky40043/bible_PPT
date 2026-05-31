"""
ppt_generator.py
直接複製 經文範本.pptx 第一頁的 shape 結構並更新文字，
保留所有原始字型大小、顏色設定（從 Slide Master 繼承）。
  Shape 0: 節號 (verse_num)
  Shape 1: 標題 (title)
  Shape 2: 經文 (body)
  Shape 3: 版本 (version)
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pathlib import Path
import copy
import io
import os

_HERE = Path(__file__).resolve().parent
_TEMPLATE_NAME = "經文範本.pptx"


def _resolve_template_path() -> str:
    env_path = os.environ.get("BIBLE_PPT_TEMPLATE")
    candidates = [
        Path(env_path) if env_path else None,
        _HERE / _TEMPLATE_NAME,
        _HERE.parent / _TEMPLATE_NAME,
        Path.cwd() / _TEMPLATE_NAME,
    ]

    for candidate in candidates:
        if candidate and candidate.exists():
            return str(candidate)

    checked = ", ".join(str(c) for c in candidates if c)
    raise FileNotFoundError(f"找不到 PPT 範本 {_TEMPLATE_NAME}; checked: {checked}")


TEMPLATE_PATH = _resolve_template_path()


def _duplicate_slide_content(ref_slide_element, dst_slide_element):
    """Replace a new slide's XML content with a deep copy of the reference slide."""
    for elem in list(dst_slide_element):
        dst_slide_element.remove(elem)
    for elem in ref_slide_element:
        dst_slide_element.append(copy.deepcopy(elem))


def _text_shape_elements(slide_element):
    """Return text-bearing shapes in slide XML order."""
    shapes = []
    for shape in slide_element.findall(".//" + qn("p:sp")):
        if shape.find(".//" + qn("p:txBody")) is not None:
            shapes.append(shape)
    return shapes


def _update_shape_text_xml(shape_element, text):
    """Update a copied shape's first text run directly in XML."""
    text_nodes = shape_element.findall(".//" + qn("a:t"))
    if not text_nodes:
        return
    text_nodes[0].text = text
    for text_node in text_nodes[1:]:
        parent = text_node.getparent()
        if parent is not None:
            parent.remove(text_node)



def generate_bible_ppt(version: str, book_zh: str, chapter: int, verses: list,
                        verse_start: int = None, verse_end: int = None,
                        include_version: bool = True) -> io.BytesIO:
    # Filter verses
    if verse_start is not None and verse_end is not None:
        filtered = [v for v in verses if verse_start <= int(v['num']) <= verse_end]
    elif verse_start is not None:
        filtered = [v for v in verses if int(v['num']) >= verse_start]
    else:
        filtered = verses

    if not filtered:
        filtered = verses

    valid_verses = [v for v in filtered if v['text'].strip()]

    # Build title label
    if verse_start is not None and verse_end is not None:
        range_label = f"{chapter}:{verse_start}-{verse_end}"
    elif verse_start is not None:
        range_label = f"{chapter}:{verse_start}"
    else:
        range_label = f"{chapter}"

    title_text = f"{book_zh} {range_label}"

    prs = Presentation(TEMPLATE_PATH)

    # Deep-copy the full reference slide XML before any modifications.
    ref_slide_element = copy.deepcopy(prs.slides[0]._element)

    # Remove all extra template slides beyond the first
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides)[1:]:
        xml_slides.remove(sldId)

    # Add slides for verses beyond the first, then duplicate slide 1's XML content.
    ref_layout = prs.slides[0].slide_layout
    for _ in valid_verses[1:]:
        slide = prs.slides.add_slide(ref_layout)
        _duplicate_slide_content(ref_slide_element, slide._element)

    # Update all slides with verse content
    for i, verse in enumerate(valid_verses):
        slide = prs.slides[i]
        tf_shapes = _text_shape_elements(slide._element)

        if len(tf_shapes) >= 1:
            _update_shape_text_xml(tf_shapes[0], str(verse['num']))
        if len(tf_shapes) >= 2:
            _update_shape_text_xml(tf_shapes[1], title_text)
        if len(tf_shapes) >= 3:
            _update_shape_text_xml(tf_shapes[2], verse['text'])
        if len(tf_shapes) >= 4:
            _update_shape_text_xml(tf_shapes[3], f"({version})" if include_version else "")

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream
